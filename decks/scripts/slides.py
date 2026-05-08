"""Reusable slide primitives for IndLokal decks."""

from __future__ import annotations

from dataclasses import dataclass

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree

from brand import (
    Theme, FONT_PRIMARY,
    SIZE_BODY_LG, SIZE_BODY, SIZE_CAPTION, SIZE_EYEBROW,
    INDIGO_900, SAFFRON_500, WHITE,
)


@dataclass
class Geometry:
    width_in: float
    height_in: float

    @property
    def margin(self): return round(self.width_in * 0.06, 3)

    @property
    def content_w(self): return self.width_in - 2 * self.margin

    @property
    def content_h(self): return self.height_in - 2 * self.margin


def _add_thick_line(slide, x1, y1, x2, y2, rgb, weight_pt):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb
    line.line.width = Pt(max(weight_pt, 4))
    line.line._get_or_add_ln().set("cap", "rnd")


def add_pulse_mark(slide, theme, x_in, y_in, size_in=0.6):
    tile = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_in), Inches(y_in), Inches(size_in), Inches(size_in),
    )
    tile.adjustments[0] = 0.22
    tile.fill.solid()
    tile.fill.fore_color.rgb = theme.brand
    tile.line.fill.background()
    tile.shadow.inherit = False
    pad = size_in * 0.18
    inner = size_in - 2 * pad
    left_x = x_in + pad
    right_x = x_in + size_in - pad
    mid_x = x_in + size_in / 2
    base_y = y_in + size_in - pad - inner * 0.18
    peak_y = y_in + pad + inner * 0.05
    _add_thick_line(slide, left_x, base_y, mid_x, peak_y, WHITE, weight_pt=size_in * 18)
    _add_thick_line(slide, mid_x, peak_y, right_x, base_y, theme.accent, weight_pt=size_in * 18)


def paint_background(slide, theme, geo):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(geo.width_in), Inches(geo.height_in))
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme.background
    bg.line.fill.background()
    bg.shadow.inherit = False


def add_pulse_motif(slide, theme, geo):
    y = geo.margin
    base = y + 0.35
    peak = y + 0.05
    seg = 0.55
    start_x = geo.width_in - geo.margin - seg * 4
    pts = [(start_x + i * seg, base if i in (0, 4) else (base - 0.05 if i in (1, 3) else peak)) for i in range(5)]
    motif_color = theme.brand if theme.name == "white" else theme.muted
    for (x1, y1), (x2, y2) in zip(pts, pts[1:]):
        line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = motif_color
        line.line.width = Pt(1.25)
        sp = line.line._get_or_add_ln()
        for child in sp.findall(qn("a:solidFill")):
            sp.remove(child)
        fill = etree.SubElement(sp, qn("a:solidFill"))
        srgb = etree.SubElement(fill, qn("a:srgbClr"), val=f"{motif_color[0]:02X}{motif_color[1]:02X}{motif_color[2]:02X}")
        etree.SubElement(srgb, qn("a:alpha"), val="22000")


def add_brand_strip(slide, theme, geo):
    add_pulse_mark(slide, theme, geo.margin, geo.margin, size_in=0.45)
    tb = slide.shapes.add_textbox(Inches(geo.margin + 0.6), Inches(geo.margin + 0.02), Inches(2.5), Inches(0.5))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "IndLokal"
    r.font.name = FONT_PRIMARY
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = theme.brand_deep if theme.name == "white" else theme.ink


def add_text(slide, text, x_in, y_in, w_in, h_in, *, size, color, bold=False,
             font=FONT_PRIMARY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = size
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_eyebrow(slide, theme, geo, text, y_in=None):
    y = y_in if y_in is not None else geo.margin + 0.7
    add_text(slide, text.upper(), geo.margin, y, geo.content_w, 0.3,
             size=SIZE_EYEBROW, color=theme.accent, bold=True)


def slide_title(prs, theme, geo, title, subtitle, eyebrow=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide, theme, geo)
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(geo.width_in * 0.62), Inches(geo.height_in),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = theme.brand_deep if theme.name == "white" else INDIGO_900
    panel.line.fill.background()
    panel.shadow.inherit = False
    mark_size = min(geo.height_in * 0.45, geo.width_in * 0.28)
    add_pulse_mark(slide, theme,
                   x_in=geo.width_in * 0.62 + (geo.width_in * 0.38 - mark_size) / 2,
                   y_in=(geo.height_in - mark_size) / 2,
                   size_in=mark_size)
    if eyebrow:
        add_text(slide, eyebrow.upper(),
                 geo.margin, geo.margin + 0.4, geo.width_in * 0.55, 0.3,
                 size=SIZE_EYEBROW, color=SAFFRON_500, bold=True)
    add_text(slide, title,
             geo.margin, geo.margin + 0.9,
             geo.width_in * 0.55, 3.0,
             size=Pt(54 if geo.width_in > 12 else 44), color=WHITE, bold=True,
             line_spacing=1.05)
    add_text(slide, subtitle,
             geo.margin, geo.margin + (3.6 if geo.width_in > 12 else 3.0),
             geo.width_in * 0.55, 2.0,
             size=Pt(20 if geo.width_in > 12 else 16),
             color=RGBColor(0xCB, 0xD5, 0xE1),
             line_spacing=1.3)
    add_text(slide, "indlokal.de  \u00b7  indlokal.com",
             geo.margin, geo.height_in - geo.margin - 0.3,
             geo.width_in * 0.55, 0.3,
             size=SIZE_CAPTION, color=RGBColor(0xA5, 0xB4, 0xFC))
    return slide


def slide_content(prs, theme, geo, eyebrow, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide, theme, geo)
    add_brand_strip(slide, theme, geo)
    add_pulse_motif(slide, theme, geo)
    if eyebrow:
        add_eyebrow(slide, theme, geo, eyebrow, y_in=geo.margin + 0.75)
    add_text(slide, title, geo.margin, geo.margin + 1.05,
             geo.content_w, 1.2,
             size=Pt(34 if geo.width_in > 12 else 28), color=theme.ink, bold=True,
             line_spacing=1.1)
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(geo.margin), Inches(geo.margin + 2.15),
        Inches(0.6), Inches(0.06),
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = theme.accent
    underline.line.fill.background()
    underline.shadow.inherit = False
    return slide, geo.margin + 2.5


def add_bullets(slide, theme, geo, items, y_in, *, size=SIZE_BODY_LG, gap=0.28):
    tb = slide.shapes.add_textbox(
        Inches(geo.margin), Inches(y_in),
        Inches(geo.content_w), Inches(geo.height_in - y_in - geo.margin - 0.4),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap * 18)
        p.line_spacing = 1.35
        r0 = p.add_run()
        r0.text = "\u258e "
        r0.font.name = FONT_PRIMARY
        r0.font.size = size
        r0.font.bold = True
        r0.font.color.rgb = theme.accent
        r = p.add_run()
        r.text = item
        r.font.name = FONT_PRIMARY
        r.font.size = size
        r.font.color.rgb = theme.ink


def add_card(slide, theme, x_in, y_in, w_in, h_in, *, title, body, icon=""):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
    )
    card.adjustments[0] = 0.08
    card.fill.solid()
    card.fill.fore_color.rgb = theme.surface
    card.line.color.rgb = theme.border
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    pad = 0.22
    cur_y = y_in + pad
    if icon:
        add_text(slide, icon, x_in + pad, cur_y, w_in - 2 * pad, 0.4,
                 size=Pt(20), color=theme.accent, bold=True)
        cur_y += 0.45
    add_text(slide, title, x_in + pad, cur_y, w_in - 2 * pad, 0.5,
             size=Pt(16), color=theme.ink, bold=True)
    cur_y += 0.45
    add_text(slide, body, x_in + pad, cur_y, w_in - 2 * pad, h_in - (cur_y - y_in) - pad,
             size=Pt(12), color=theme.muted, line_spacing=1.35)


def add_stat(slide, theme, x_in, y_in, w_in, h_in, *, big, label, sub=""):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
    )
    card.adjustments[0] = 0.08
    card.fill.solid()
    card.fill.fore_color.rgb = theme.brand if theme.name == "white" else theme.surface
    card.line.fill.background()
    card.shadow.inherit = False
    pad = 0.22
    big_color = WHITE if theme.name == "white" else theme.accent
    label_color = RGBColor(0xE0, 0xE7, 0xFF) if theme.name == "white" else theme.ink
    sub_color = RGBColor(0xC7, 0xD2, 0xFE) if theme.name == "white" else theme.muted
    add_text(slide, big, x_in + pad, y_in + pad, w_in - 2 * pad, h_in * 0.5,
             size=Pt(36), color=big_color, bold=True)
    add_text(slide, label, x_in + pad, y_in + pad + h_in * 0.48, w_in - 2 * pad, 0.4,
             size=Pt(13), color=label_color, bold=True)
    if sub:
        add_text(slide, sub, x_in + pad, y_in + pad + h_in * 0.7, w_in - 2 * pad, h_in * 0.3,
                 size=Pt(10), color=sub_color, line_spacing=1.3)


def slide_section(prs, theme, geo, label, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide, theme, geo)
    bg_color = theme.brand_deep if theme.name == "white" else theme.surface
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(geo.width_in), Inches(geo.height_in)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = bg_color
    panel.line.fill.background()
    panel.shadow.inherit = False
    add_pulse_mark(slide, theme, geo.width_in - geo.margin - 0.6, geo.margin, size_in=0.6)
    add_text(slide, label.upper(),
             geo.margin, geo.height_in / 2 - 0.9, geo.content_w, 0.4,
             size=SIZE_EYEBROW, color=SAFFRON_500, bold=True)
    add_text(slide, title,
             geo.margin, geo.height_in / 2 - 0.4, geo.content_w, 2.0,
             size=Pt(46 if geo.width_in > 12 else 36), color=WHITE, bold=True,
             line_spacing=1.1)
    return slide


def slide_closing(prs, theme, geo, headline, sub, contact_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide, theme, geo)
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(geo.width_in), Inches(geo.height_in)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = theme.brand_deep if theme.name == "white" else theme.surface
    panel.line.fill.background()
    panel.shadow.inherit = False
    add_pulse_mark(slide, theme, geo.margin, geo.margin, size_in=0.6)
    add_text(slide, "IndLokal",
             geo.margin + 0.75, geo.margin + 0.07, 3.0, 0.5,
             size=Pt(22), color=WHITE, bold=True)
    add_text(slide, headline,
             geo.margin, geo.margin + 1.6, geo.content_w, 2.0,
             size=Pt(48 if geo.width_in > 12 else 36), color=WHITE, bold=True,
             line_spacing=1.1)
    add_text(slide, sub,
             geo.margin, geo.margin + (3.4 if geo.width_in > 12 else 2.8),
             geo.content_w, 1.2,
             size=Pt(20 if geo.width_in > 12 else 16),
             color=RGBColor(0xCB, 0xD5, 0xE1), line_spacing=1.35)
    y = geo.height_in - geo.margin - 1.2
    for line in contact_lines:
        add_text(slide, line, geo.margin, y, geo.content_w, 0.35,
                 size=Pt(14), color=RGBColor(0xE0, 0xE7, 0xFF))
        y += 0.32


def new_presentation(width_in, height_in):
    prs = Presentation()
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)
    return prs, Geometry(width_in, height_in)
